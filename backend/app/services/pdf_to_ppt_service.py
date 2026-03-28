import os
import io
from typing import List, Optional
from fastapi import UploadFile, HTTPException
from pdf2image import convert_from_path
from app.utils.file_handler import save_upload_files, create_temp_file


async def run(files: List[UploadFile], options: Optional[dict] = None) -> str:
    paths: List[str] = []
    try:
        paths = await save_upload_files(files, allowed_exts=[".pdf"])
        pages = convert_from_path(paths[0])
        if not pages:
            raise HTTPException(status_code=400, detail="No pages in PDF")
        try:
            from pptx import Presentation
            from pptx.util import Inches
            out_path, out_id = create_temp_file(".pptx")
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            for img in pages:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                b = io.BytesIO()
                img.save(b, format="JPEG")
                b.seek(0)
                slide.shapes.add_picture(b, 0, 0, prs.slide_width, prs.slide_height)
            prs.save(out_path)
            return out_id
        except Exception:
            from app.services.convert_service import pdf_to_ppt as fallback
            out_id = fallback(paths[0])
            return out_id
    finally:
        for p in paths:
            try:
                os.remove(p)
            except Exception:
                pass
